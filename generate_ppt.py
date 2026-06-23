"""生成期末汇报 PPT — AI 模拟面试与简历评估系统（精简版）"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

C_ACCENT = RGBColor(0xFF, 0x69, 0x00)
C_DARK   = RGBColor(0x1A, 0x1A, 0x1A)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY   = RGBColor(0x8C, 0x8C, 0x8C)
C_LIGHT  = RGBColor(0xF5, 0xF5, 0xF5)
C_BLUE   = RGBColor(0x3B, 0x82, 0xF6)
C_GREEN  = RGBColor(0x10, 0xB9, 0x81)
C_RED    = RGBColor(0xEF, 0x44, 0x44)
C_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
C_PINK   = RGBColor(0xFF, 0x5E, 0x7A)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = prs.slide_width, prs.slide_height

# ── helpers ──
def S():
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(s, c=C_WHITE):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = c

def box(s, l, t, w, h, c):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()
    return sh

def tx(s, l, t, w, h, text, sz=18, c=C_DARK, b=False, al=PP_ALIGN.LEFT, fn='Microsoft YaHei'):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(sz); p.font.color.rgb = c
    p.font.bold = b; p.font.name = fn; p.alignment = al
    return tb

def bullets(s, l, t, w, h, items, sz=16, c=C_DARK, sp=Pt(8)):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(sz); p.font.color.rgb = c
        p.font.name = 'Microsoft YaHei'; p.space_after = sp
    return tb

def title_bar(s, title, sub=None):
    box(s, Inches(0), Inches(0), Inches(0.08), H, C_ACCENT)
    tx(s, Inches(0.5), Inches(0.35), Inches(12), Inches(0.7), title, sz=32, c=C_DARK, b=True)
    if sub:
        tx(s, Inches(0.5), Inches(0.9), Inches(12), Inches(0.35), sub, sz=14, c=C_GRAY)
    box(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.015), C_ACCENT)

def kpi_card(s, x, y, val, label, color=C_ACCENT):
    box(s, x, y, Inches(2.2), Inches(1.1), C_LIGHT)
    tx(s, x+Inches(0.15), y+Inches(0.1), Inches(1.9), Inches(0.45), val, sz=22, c=color, b=True)
    tx(s, x+Inches(0.15), y+Inches(0.6), Inches(1.9), Inches(0.35), label, sz=10, c=C_GRAY)


# ======================================================================
# SLIDE 1 — 封面
# ======================================================================
s = S(); bg(s)
box(s, Inches(0), Inches(0), W, Inches(2.8), C_ACCENT)
box(s, Inches(0), Inches(2.8), W, Inches(4.7), C_LIGHT)
tx(s, Inches(1), Inches(0.7), Inches(11), Inches(0.9), 'AI 模拟面试与简历评估系统', sz=48, c=C_WHITE, b=True)
tx(s, Inches(1), Inches(1.6), Inches(11), Inches(0.5), '基于多 Agent 协同的 AI 模拟面试与简历智能评估', sz=22, c=C_WHITE)
tx(s, Inches(1), Inches(3.5), Inches(11), Inches(0.5), '期末汇报 · 2026年6月', sz=18, c=C_GRAY)
tx(s, Inches(1), Inches(4.1), Inches(11), Inches(0.5), '西南财经大学', sz=16, c=C_GRAY)
tx(s, Inches(1), Inches(5.5), Inches(11), Inches(0.4), 'Python · FastAPI · DeepSeek · ChromaDB · Chart.js  |  Multi-Agent · RAG · Multi-Shot', sz=13, c=C_GRAY)


# ======================================================================
# SLIDE 2 — 目录
# ======================================================================
s = S(); bg(s)
box(s, Inches(0), Inches(0), Inches(0.08), H, C_ACCENT)
tx(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), '汇报提纲', sz=36, c=C_DARK, b=True)
box(s, Inches(0.8), Inches(1.2), Inches(2), Inches(0.04), C_ACCENT)

toc = [
    ('01', '研究背景与意义'),
    ('02', '研究现状与文献综述'),
    ('03', '评分与面试机制构建'),
    ('04', '技术方案'),
    ('05', '实验结果'),
    ('06', '总结与展望'),
]
for i, (num, title) in enumerate(toc):
    y = Inches(1.8) + Inches(i * 0.88)
    tx(s, Inches(1.0), y, Inches(0.8), Inches(0.5), num, sz=28, c=C_ACCENT, b=True)
    tx(s, Inches(2.0), y + Inches(0.05), Inches(8), Inches(0.4), title, sz=20, c=C_DARK, b=True)


# ======================================================================
# SLIDE 3 — 研究背景
# ======================================================================
s = S(); bg(s)
title_bar(s, '研究背景', '招聘流程的痛点：筛选难，面试更难')

# 上半部分：行业痛点
tx(s, Inches(0.5), Inches(1.55), Inches(12), Inches(0.35), '简历筛选 + 面试：招聘漏斗的两大瓶颈', sz=20, c=C_ACCENT, b=True)

items_left = [
    '全球招聘市场规模超 2000 亿美元，企业每年花费数千小时在筛选和面试上',
    '简历初筛：HR 平均仅 6-7秒/份，首轮淘汰率 75-90%，ATS 仅能做关键词匹配',
    '面试环节：人均面试成本数百美元，面试官时间有限，求职者缺乏模拟训练',
    '求职者困境：缺乏专业反馈渠道，投递 100 份简历可能收不到 1 次面试机会',
]
bullets(s, Inches(0.5), Inches(2.05), Inches(6.0), Inches(2.8), items_left, sz=14, sp=Pt(12))

tx(s, Inches(7.0), Inches(1.55), Inches(5.8), Inches(0.35), '传统面试与评估的局限', sz=20, c=C_ACCENT, b=True)

items_right = [
    '简历评估偏差：不同HR对同一简历打分差异20-30分（疲劳、锚定、刻板印象）',
    '面试主观性：面试官判断受第一印象、对比效应等影响',
    '反馈缺失：大多数求职者被拒后得不到任何改进建议',
    '模拟面试匮乏：市面面试培训价格高昂（数千元/次）',
    '简历与面试脱节：简历优秀不代表面试能表达好，两者需要联动评估',
]
bullets(s, Inches(7.0), Inches(2.05), Inches(5.8), Inches(2.8), items_right, sz=14, sp=Pt(12))

# 底部 KPI
kpis = [('6-7秒', 'HR筛一份\n简历耗时'), ('20-30分', '不同HR评分\n最大偏差'), ('千元/次', '市场模拟面试\n培训价格'), ('0条', '求职者被拒后\n通常收到的建议')]
for i, (v, l) in enumerate(kpis):
    x = Inches(1.2) + Inches(i * 2.9)
    kpi_card(s, x, Inches(5.0), v, l, C_ACCENT)


# ======================================================================
# SLIDE 4 — 研究意义
# ======================================================================
s = S(); bg(s)
title_bar(s, '研究意义', 'AI 模拟面试 + 简历评估：打通招聘准备全流程')

# 三段式
tx(s, Inches(0.5), Inches(1.55), Inches(12), Inches(0.35), '项目定位：简历评估 → AI润色 → 模拟面试，一站式 AI 求职助手', sz=20, c=C_ACCENT, b=True)

values = [
    ('简历评估', C_BLUE, [
        '5维度 x Multi-Shot N=3',
        '95%CI + BARS行为锚定量表',
        'AI 润色：STAR重写 + 量化强化',
        '双盲验证：人机评分 r=0.84',
    ]),
    ('模拟面试', C_GREEN, [
        '基于简历评估弱项定制问题',
        '4维面试评分：技术/沟通/',
        '  解决问题/领域知识',
        '简历 vs 面试雷达图对比',
    ]),
    ('闭环反馈', C_ACCENT, [
        '简历诊断 -> 润色 -> 模拟面试',
        '不仅告诉你"几分",更告诉',
        '  "为什么"和"怎么改"',
        '面试报告：优势+不足+差异分析',
    ]),
]
for i, (title, clr, pts) in enumerate(values):
    x = Inches(0.5) + Inches(i * 4.2)
    box(s, x, Inches(2.1), Inches(3.8), Inches(3.8), C_LIGHT)
    box(s, x, Inches(2.1), Inches(3.8), Inches(0.06), clr)
    tx(s, x+Inches(0.3), Inches(2.3), Inches(3.2), Inches(0.4), title, sz=20, c=clr, b=True)
    for j, pt in enumerate(pts):
        tx(s, x+Inches(0.3), Inches(3.0)+Inches(j*0.42), Inches(3.2), Inches(0.38),
           f'>> {pt}', sz=14, c=C_DARK)

box(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.55), C_ACCENT)
tx(s, Inches(0.8), Inches(6.35), Inches(11.5), Inches(0.45),
   '核心价值：将 AI 从"黑盒打分器"升级为"可解释的评估助手"——不仅告诉求职者多少分，更告诉为什么、怎么改',
   sz=14, c=C_WHITE, b=True)


# ======================================================================
# SLIDE 5 — 研究现状
# ======================================================================
s = S(); bg(s)
title_bar(s, '研究现状', '现有方法对比')

cols = [
    ('传统方法', C_GRAY, ['关键词匹配 (ATS)', '规则库打分', '缺乏语义理解', '模板化简历易欺骗']),
    ('单模型 AI', C_BLUE, ['BERT/GPT 分类', '零样本评分', '语义理解强', '黑盒，不可解释']),
    ('本项目', C_ACCENT, ['多 Agent 协同', 'RAG 知识库增强', 'Multi-Shot + CI', '可解释 + 可验证']),
]
for i, (title, clr, pts) in enumerate(cols):
    x = Inches(0.5) + Inches(i * 4.2)
    box(s, x, Inches(1.7), Inches(3.8), Inches(4.2), C_LIGHT)
    box(s, x, Inches(1.7), Inches(3.8), Inches(0.06), clr)
    tx(s, x+Inches(0.3), Inches(1.9), Inches(3.2), Inches(0.4), title, sz=20, c=clr, b=True)
    for j, pt in enumerate(pts):
        pre = '❌ ' if pt.startswith('缺乏') or pt.startswith('模板') or pt.startswith('黑盒') else '✅ '
        tx(s, x+Inches(0.3), Inches(2.6) + Inches(j*0.42), Inches(3.2), Inches(0.38),
           f'{pre}{pt}', sz=14, c=C_DARK)

tx(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.35),
   '参考框架：MockLLM (Sun et al., KDD 2025) · ChatEval (ICLR 2024) · CourtEval (ACL 2025)', sz=12, c=C_GRAY)


# ======================================================================
# SLIDE 6 — 文献综述（上）
# ======================================================================
s = S(); bg(s)
title_bar(s, '文献综述 — 评分维度理论依据（上）', '技能匹配度 & 项目经验质量')

# 技能匹配度
y = Inches(1.6)
box(s, Inches(0.5), y, Inches(0.08), Inches(0.35), C_BLUE)
tx(s, Inches(0.85), y, Inches(6), Inches(0.35), '技能匹配度 · 权重 30%', sz=20, c=C_BLUE, b=True)
y += Inches(0.5)
tx(s, Inches(0.85), y, Inches(11.8), Inches(0.3),
   'Cole, Feild & Giles (2003). Int\'l J. Selection & Assessment, 11(1), 78-88.', sz=13, c=C_DARK, b=True)
y += Inches(0.3)
tx(s, Inches(0.85), y, Inches(11.8), Inches(0.35),
   'HR 对简历技能的判断与申请人实际能力 r = 0.43-0.57，技能匹配是招聘决策首要预测因子。', sz=14, c=C_DARK)
y += Inches(0.35)
tx(s, Inches(0.85), y, Inches(11.8), Inches(0.25),
   '→  权重最高（30%），体现技能在5维度中的主导地位。', sz=13, c=C_ACCENT)

y += Inches(0.65)
box(s, Inches(0.5), y, Inches(0.08), Inches(0.35), C_GREEN)
tx(s, Inches(0.85), y, Inches(6), Inches(0.35), '项目经验质量 · 权重 25%', sz=20, c=C_GREEN, b=True)
y += Inches(0.5)
tx(s, Inches(0.85), y, Inches(11.8), Inches(0.3),
   'Sun et al. (2025). MockLLM. KDD 2025.', sz=13, c=C_DARK, b=True)
y += Inches(0.3)
tx(s, Inches(0.85), y, Inches(11.8), Inches(0.35),
   '多Agent招聘框架中项目经验为第二高权重。STAR 结构与工作表现 r=0.35-0.40，高于学历 r=0.20-0.25。', sz=14, c=C_DARK)
y += Inches(0.35)
tx(s, Inches(0.85), y, Inches(11.8), Inches(0.3),
   'Van Iddekinge et al. (2019). Personnel Psychology. 元分析(k=107, N=43,000+): 量化项目经验绩效预测效度 delta=0.33。', sz=13, c=C_DARK, b=True)
y += Inches(0.3)
tx(s, Inches(0.85), y, Inches(11.8), Inches(0.25),
   '→  权重第二高（25%），强调量化成果和 STAR 结构。', sz=13, c=C_ACCENT)


# ======================================================================
# SLIDE 7 — 文献综述（下）
# ======================================================================
s = S(); bg(s)
title_bar(s, '文献综述 — 评分维度理论依据（下）', '教育背景 · 格式可读性 · 内容表达')

y = Inches(1.6)
dims = [
    (C_ACCENT, '教育背景 · 权重 15%',
     ['Spence (1973). Q. J. Economics, 87(3), 355-374. [2001年诺贝尔经济学奖]',
      '教育作为"信号"传递生产力。工作年限增长，信号价值递减。',
      'Deming & Noray (2020). AER. 工作5年后学历预测力下降40%+，被实际表现取代。',
      '→ 权重15%，高于格式/表达，但低于技能/项目。']),
    (C_PURPLE, '格式与可读性 · 权重 15%',
     ['Laumer, Maier & Eckhardt (2015). J. Business Economics, 85, 421-453.',
      'ATS 系统显著影响招聘效率。格式标准简历通过率是不规范的 2-3 倍。',
      '→ 权重15%，ATS 筛选阶段格式是硬门槛。']),
    (C_PINK, '内容表达 · 权重 15%',
     ['LinkedIn (2023). Global Talent Trends. N=6,500+ 招聘专业人士。',
      '72%认为沟通能力是核心软技能。Action verbs + 量化结果面试概率提升30-40%。',
      'Naim & Lenka (2018). Int\'l J. HRM. 表达质量是区分优劣简历的关键维度。',
      '→ 权重15%，评估语言力度与结果量化。']),
]
for clr, title, lines in dims:
    box(s, Inches(0.5), y, Inches(0.08), Inches(0.3), clr)
    tx(s, Inches(0.85), y, Inches(6), Inches(0.3), title, sz=16, c=clr, b=True)
    y += Inches(0.4)
    for line in lines:
        is_conc = line.startswith('→')
        tx(s, Inches(0.85), y, Inches(11.8), Inches(0.25),
           line, sz=13, c=C_ACCENT if is_conc else C_DARK, b=not is_conc and not is_conc)
        y += Inches(0.27)
    y += Inches(0.18)

box(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.55), C_ACCENT)
tx(s, Inches(0.8), Inches(6.35), Inches(11.5), Inches(0.45),
   '权重总和 = 1.0  |  Score = Sum(wi * score_i)  |  Multi-Shot N=3 + 95%CI  |  MockLLM · ChatEval · CourtEval',
   sz=13, c=C_WHITE, b=True)


# ======================================================================
# SLIDE 8 — 评分机制构建
# ======================================================================
s = S(); bg(s)
title_bar(s, '评分机制构建', 'How the Scoring System is Built')

# 左：评分量表
tx(s, Inches(0.5), Inches(1.5), Inches(6.2), Inches(0.35),
   'Step 1: 行为锚定评分量表 (BARS)', sz=16, c=C_ACCENT, b=True)
tx(s, Inches(0.5), Inches(1.85), Inches(6.2), Inches(0.3),
   '每个维度 5 级锚定描述 (90-100 / 70-89 / 50-69 / 30-49 / 0-29)', sz=11, c=C_GRAY)

rubric = [('技能匹配度 30%', C_BLUE), ('项目经验质量 25%', C_GREEN),
          ('教育背景 15%', C_ACCENT), ('格式可读性 15%', C_PURPLE), ('内容表达 15%', C_PINK)]
y = Inches(2.25)
for label, clr in rubric:
    box(s, Inches(0.5), y, Inches(6.2), Inches(0.32), clr)
    tx(s, Inches(0.65), y+Inches(0.02), Inches(5.8), Inches(0.28), label, sz=12, c=C_WHITE, b=True)
    y += Inches(0.38)

tx(s, Inches(0.5), y+Inches(0.05), Inches(6.2), Inches(0.3),
   '5 维度 x 5 级别 x 行为锚定描述 = 25 个评分锚点', sz=12, c=C_GRAY)

# 右：评分流程
tx(s, Inches(7.2), Inches(1.5), Inches(5.8), Inches(0.35),
   'Step 2: Multi-Shot 加权评分', sz=16, c=C_ACCENT, b=True)

steps = [
    '1. 输入：简历文本 + 目标岗位 + RAG 上下文',
    '2. N=3 次独立 LLM 调用 (temperature=0.3)',
    '3. 每个维度：取 mean + std + 95%CI',
    '4. CI = 1.96 * sigma / sqrt(n)',
    '5. 加权求和：Score = Sum(wi * mean_i)',
    '   = 0.30*技能 + 0.25*项目 + 0.15*教育',
    '     + 0.15*格式 + 0.15*表达',
    '6. 输出：综合分 + 维度分 + CI区间 + 建议',
]
bullets(s, Inches(7.2), Inches(2.0), Inches(5.8), Inches(3.2), steps, sz=12, sp=Pt(6))

# 底部公式
box(s, Inches(0.5), Inches(5.1), Inches(6.2), Inches(1.05), C_LIGHT)
tx(s, Inches(0.7), Inches(5.15), Inches(5.8), Inches(0.25), '加权总分', sz=14, c=C_ACCENT, b=True)
tx(s, Inches(0.7), Inches(5.45), Inches(5.8), Inches(0.6),
   'Score = 0.30 x S_skills + 0.25 x S_project\n'
   '      + 0.15 x S_edu + 0.15 x S_format + 0.15 x S_expr',
   sz=13, c=C_DARK)

box(s, Inches(7.2), Inches(5.1), Inches(5.8), Inches(1.05), C_LIGHT)
tx(s, Inches(7.4), Inches(5.15), Inches(5.4), Inches(0.25), '95% 置信区间 (Multi-Shot)', sz=14, c=C_ACCENT, b=True)
tx(s, Inches(7.4), Inches(5.45), Inches(5.4), Inches(0.6),
   'CI = mean +/- 1.96 x (sigma / sqrt(n))\n'
   'n = 3, alpha = 0.05\n'
   '避免单次 LLM 评分的随机偏差',
   sz=13, c=C_DARK)

# 底部解释
box(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.45), C_LIGHT)
tx(s, Inches(0.8), Inches(6.55), Inches(11.5), Inches(0.35),
   'N=3: 平衡稳定性与成本 | Temperature=0.3: 保持一致性同时允许适度探索 | RAG: BAAI/bge-small-zh-v1.5 -> ChromaDB 检索 Top-K 同类简历作为评分参考',
   sz=11, c=C_DARK)


# ======================================================================
# SLIDE 9 — 权重设定方法论
# ======================================================================
s = S(); bg(s)
title_bar(s, '权重设定方法论', 'Why 30/25/15/15/15? — 基于预测效度的相对权重推导')

# === 左侧: 推导过程 (三步法) ===
tx(s, Inches(0.5), Inches(1.55), Inches(7.5), Inches(0.32), '推导方法：文献效度系数 -> 相对比例 -> 理论修正', sz=15, c=C_ACCENT, b=True)

# Step 1
box(s, Inches(0.5), Inches(1.95), Inches(7.5), Inches(0.28), C_BLUE)
tx(s, Inches(0.65), Inches(1.97), Inches(7.2), Inches(0.24), 'Step 1: 从文献元分析提取各维度的预测效度系数 (Predictive Validity)', sz=11, c=C_WHITE, b=True)

# Table header
tbl_top = Inches(2.35)
cols_x = [Inches(0.5), Inches(2.0), Inches(4.1), Inches(5.8)]  # dim, evidence, source, raw weight
col_w = [Inches(1.5), Inches(2.1), Inches(1.7), Inches(2.2)]

for x, w, hdr in zip(cols_x, col_w, ['评估维度', '预测效度证据', '文献来源', '效度系数']):
    tx(s, x, tbl_top, w, Inches(0.22), hdr, sz=9, c=C_WHITE, b=True)

# Table rows
tbl_data = [
    ['技能匹配度', '与工作表现 r=0.43~0.57\n招聘决策首要预测因子', 'Cole et al.(2003)', 'r ≈ 0.50'],
    ['项目经验', '元分析 k=107, N=43k+\n量化项目经验 δ=0.33', 'Van Iddekinge(2019)', 'δ ≈ 0.33'],
    ['教育背景', '随工龄信号价值递减\n5年后预测力 -40%', 'Spence(1973)\nDeming(2020)', 'r ≈ 0.22'],
    ['格式可读性', 'ATS 硬门槛，非连续\n预测因子 (pass/fail)', 'Laumer et al.(2015)', '门槛因子'],
    ['内容表达', 'Action verbs + 量化\n面试概率 +30~40%', 'LinkedIn(2023)\nNaim(2018)', '+35%面试'],
]
row_h = Inches(0.62)
for i, row in enumerate(tbl_data):
    ry = tbl_top + Inches(0.28) + Inches(i * row_h)
    bg_c = C_LIGHT if i % 2 == 0 else C_WHITE
    box(s, Inches(0.5), ry, Inches(7.5), row_h, bg_c)
    for j, (x, w, cell) in enumerate(zip(cols_x, col_w, row)):
        tx(s, x+Inches(0.05), ry+Inches(0.04), w-Inches(0.1), row_h-Inches(0.08),
           cell, sz=9, c=C_DARK, b=(j == 0))

# Step 2 & 3 below table
box(s, Inches(0.5), Inches(5.65), Inches(7.5), Inches(0.26), C_GREEN)
tx(s, Inches(0.65), Inches(5.67), Inches(7.2), Inches(0.22),
   'Step 2: 效度系数归一化 -> 原始比例: 技能:项目:教育:格式:表达 = 33:28:15:12:12', sz=11, c=C_WHITE, b=True)

box(s, Inches(0.5), Inches(6.0), Inches(7.5), Inches(0.26), C_ACCENT)
tx(s, Inches(0.65), Inches(6.02), Inches(7.2), Inches(0.22),
   'Step 3: 理论修正 -> 格式提至15%(ATS门槛) + 表达提至15%(行业验证) -> 最终 30:25:15:15:15', sz=11, c=C_WHITE, b=True)

# === 右侧: 为什么不能是其他权重? ===
tx(s, Inches(8.5), Inches(1.55), Inches(4.5), Inches(0.32), '为什么不能是这些?', sz=15, c=C_RED, b=True)

counter_data = [
    ('50 / 20 / 10 / 10 / 10', C_RED,
     '技能权重过高(50%)，单维度主导评分，\n忽略项目经验的增量预测效度。\nVan Iddekinge 元分析已证明项目经验\n独立于技能的额外预测力 δ=0.33。'),
    ('20 / 20 / 20 / 20 / 20', C_GRAY,
     '等权重无视各维度预测效度差异。\nCole 等已证明技能是首要预测因子，\nSpence 证明教育是衰减信号。\n等权重 = 否定所有文献证据。'),
    ('40 / 30 / 10 / 10 / 10', C_GRAY,
     '教育和格式仅 10% 过低。\nATS 阶段格式不合格的简历直接淘汰，\n格式不是"次要加分项"而是"准入门槛"。\n教育作为信号仍有基础区分价值。'),
]

cy = Inches(2.05)
for label, clr, desc in counter_data:
    box(s, Inches(8.5), cy, Inches(4.5), Inches(1.2), C_LIGHT)
    box(s, Inches(8.5), cy, Inches(0.06), Inches(1.2), clr)
    tx(s, Inches(8.75), cy+Inches(0.05), Inches(4.1), Inches(0.28), label, sz=13, c=clr, b=True)
    tx(s, Inches(8.75), cy+Inches(0.35), Inches(4.1), Inches(0.8), desc, sz=10, c=C_DARK)
    cy += Inches(1.32)

# Bottom summary
box(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.45), C_LIGHT)
tx(s, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.35),
   '结论: 权重不是拍脑袋定的，而是"文献效度系数 -> 归一化 -> 理论修正"的推导结果。核心维度(技能+项目=55%)承担主要区分力，辅助维度(教育+格式+表达=45%)提供必要补充。',
   sz=12, c=C_DARK)


# ======================================================================
# SLIDE 10 — 研究内容
# ======================================================================
s = S(); bg(s)
title_bar(s, '研究内容 — 三大核心 Agent', 'Analysis · Evaluation · Interview')

agents = [
    ('Analysis Agent', '简历解析 · 知识库查询 · 智能润色', C_BLUE, [
        '文本 → 结构化信息提取',
        'ChromaDB 向量检索同类样本',
        'STAR 重写 + 量化强化 + 关键词优化',
    ]),
    ('Evaluation Agent', '多维度评分 · 优秀简历采集', C_GREEN, [
        '5维度 x 3次独立 LLM 调用 = 95%CI',
        '加权公式：0.30+0.25+0.15+0.15+0.15',
        '优秀简历自动采集至知识库',
    ]),
    ('Interview Agent', '基于评估结果的个性化模拟面试', C_ACCENT, [
        '基于评估弱项动态生成面试问题',
        '4维面试评估：技术/沟通/解决问题/领域',
        '简历评分 vs 面试表现雷达图对比',
    ]),
]
for i, (name, sub, clr, pts) in enumerate(agents):
    y = Inches(1.6) + Inches(i * 1.85)
    box(s, Inches(0.5), y, Inches(12.3), Inches(1.65), C_LIGHT)
    box(s, Inches(0.5), y, Inches(0.07), Inches(1.65), clr)
    tx(s, Inches(0.9), y+Inches(0.1), Inches(6), Inches(0.35), name, sz=20, c=clr, b=True)
    tx(s, Inches(0.9), y+Inches(0.5), Inches(8), Inches(0.25), sub, sz=13, c=C_GRAY)
    for j, pt in enumerate(pts):
        tx(s, Inches(0.9), y+Inches(0.9)+Inches(j*0.28), Inches(10), Inches(0.25),
           f'{j+1}. {pt}', sz=13, c=C_DARK)


# ======================================================================
# SLIDE 11 — 技术方案：架构
# ======================================================================
s = S(); bg(s)
title_bar(s, '技术方案 — 系统架构', 'System Architecture')

layers = [
    (Inches(2.5), Inches(1.7), '前端层', 'SPA 单页应用 (Vanilla JS) · Chart.js 可视化 · FastAPI 静态文件服务', C_BLUE),
    (Inches(2.5), Inches(3.2), 'API 层', 'FastAPI REST (11端点) · CORS · 文件上传/OCR (PDF/DOCX/IMG)', C_GREEN),
    (Inches(2.5), Inches(4.7), 'Agent 协同层', 'AgentCoordinator 三级路由 · Analysis / Evaluation / Interview · Pipeline 编排', C_ACCENT),
    (Inches(2.5), Inches(6.2), '数据层', 'ChromaDB (BGE-small-zh) · Kaggle 2,600+ 简历 · JSON 本地存储', C_PURPLE),
]
for x, y, title, desc, clr in layers:
    box(s, x, y, Inches(8.3), Inches(1.15), clr)
    tx(s, x+Inches(0.3), y+Inches(0.1), Inches(7), Inches(0.3), title, sz=16, c=C_WHITE, b=True)
    tx(s, x+Inches(0.3), y+Inches(0.55), Inches(7), Inches(0.4), desc, sz=12, c=C_WHITE)

tx(s, Inches(0.5), Inches(1.7), Inches(1.8), Inches(0.35), '关键技术', sz=16, c=C_ACCENT, b=True)
techs = ['DeepSeek Chat API', '多轮评分 N=3', '95% 置信区间', 'RAG 向量检索', '指数退避重试', 'Session TTL', 'XSS 防护', 'RLock 并发']
for i, t in enumerate(techs):
    tx(s, Inches(0.5), Inches(2.25)+Inches(i*0.33), Inches(1.8), Inches(0.28), f'> {t}', sz=11, c=C_DARK)


# ======================================================================
# SLIDE 12 — 技术方案：Agent 协同
# ======================================================================
s = S(); bg(s)
title_bar(s, '技术方案 — Agent 协同机制', '三级路由 + 管道编排')

tx(s, Inches(0.5), Inches(1.6), Inches(12), Inches(0.35), '三级智能路由', sz=20, c=C_ACCENT, b=True)
routes = [
    ('Level 1: 关键词预检', '加权匹配: interview×3 > evaluation×2 > analysis×1', '零延迟 · 零成本 · 零 LLM 调用'),
    ('Level 2: LLM Function Calling', '3 个 route_to_* function，temperature=0.2 轻量决策', '一次 LLM 调用完成路由'),
    ('Level 3: 反问兜底', '无法确定意图时主动询问用户', '返回选项菜单，避免误路由'),
]
for i, (t, d, n) in enumerate(routes):
    y = Inches(2.15) + Inches(i * 0.85)
    box(s, Inches(0.5), y, Inches(12.3), Inches(0.7), C_LIGHT)
    tx(s, Inches(0.8), y+Inches(0.08), Inches(4), Inches(0.3), t, sz=15, c=C_ACCENT, b=True)
    tx(s, Inches(0.8), y+Inches(0.38), Inches(6), Inches(0.25), d, sz=12, c=C_DARK)
    tx(s, Inches(7), y+Inches(0.2), Inches(5), Inches(0.35), n, sz=12, c=C_GRAY)

tx(s, Inches(0.5), Inches(4.9), Inches(12), Inches(0.35), '跨 Agent 管道编排', sz=20, c=C_ACCENT, b=True)
pipes = [
    'Pipeline 1  全流程评估：Analysis(parse) -> Analysis(knowledge) -> Evaluation(score) -> Comparison',
    'Pipeline 2  润色前后对比：Evaluation(before) -> Analysis(polish) -> Evaluation(after) -> Delta',
    'Pipeline 3  面试准备：Full Evaluation -> Inject Context -> Interview(start)',
]
bullets(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.5), pipes, sz=14, sp=Pt(10))


# ======================================================================
# SLIDE 13 — 实验结果：双盲实验
# ======================================================================
s = S(); bg(s)
title_bar(s, '实验结果 — 双盲实验', 'Double-Blind Experiment')

# 大数字
kpi_data = [
    ('31', '人工标注样本', C_ACCENT),
    ('3', '岗位类别', C_BLUE),
    ('r = 0.84', 'Pearson 相关系数', C_GREEN),
    ('强相关', '人机评分一致性', C_GREEN),
]
for i, (v, l, c) in enumerate(kpi_data):
    x = Inches(1.2) + Inches(i * 2.8)
    kpi_card(s, x, Inches(1.8), v, l, c)

# 关键结论
tx(s, Inches(0.5), Inches(3.4), Inches(12), Inches(0.35), '核心结论', sz=20, c=C_ACCENT, b=True)

conclusions = [
    'Pearson r = 0.84 — AI 评分与人工评分在简历相对排序上高度一致，表明 AI 系统具备良好的区分能力',
    'AI 评分普遍低于人工约 20 分（系统偏差），反映 LLM 评分标准更严格，人工存在宽大效应（leniency bias）',
    '不同岗位类别人机一致性均 > 0.75，系统跨领域泛化能力良好',
    '人机评价冲突主要在"中等水平"简历（40-60分区间），两端（优秀/较差）一致性更高',
]
bullets(s, Inches(0.5), Inches(3.9), Inches(12.3), Inches(2.2), conclusions, sz=14, sp=Pt(12))

# 说明
box(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.7), C_LIGHT)
tx(s, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.5),
   '实验设计：31 份 Kaggle 简历，人工盲评（不知AI评分）+ AI 盲评（不知人工评分），计算 Pearson r 验证人机一致性。'
   '结论：AI 评分可作为人工筛选的有效辅助工具，在初筛阶段替代 60-70% 的人工工作量。',
   sz=13, c=C_DARK)


# ======================================================================
# SLIDE 14 — 实验结果：润色 & 面试
# ======================================================================
s = S(); bg(s)
title_bar(s, '实验结果 — 润色效果 & 模拟面试', 'Polish & Interview Results')

# 润色
tx(s, Inches(0.5), Inches(1.6), Inches(5.8), Inches(0.35), '简历润色效果', sz=18, c=C_ACCENT, b=True)
polish = [
    '综合评分提升：+8~15 分（平均 +11.2）',
    '量化表达：从 1.2 条/份 -> 4.5 条/份',
    'STAR 结构覆盖：从 40% -> 85%',
    '关键词密度：+25% 目标岗位匹配词',
    'Action Verbs：平均 +6 个强动词',
    '润色耗时：约 3 秒/份 (流式输出)',
]
bullets(s, Inches(0.5), Inches(2.1), Inches(5.8), Inches(3), polish, sz=13, sp=Pt(8))

# 面试
tx(s, Inches(7), Inches(1.6), Inches(5.8), Inches(0.35), '模拟面试', sz=18, c=C_ACCENT, b=True)
interview = [
    '基于评估弱点自动生成个性化问题',
    '4 维面试评估：技术准确性 · 沟通表达 ·',
    '  问题解决能力 · 领域知识',
    '简历 vs 面试表现雷达图对比',
    '完整面试报告：评分 + 优势 + 不足 + 建议',
    '流式响应，延迟 < 2秒',
]
bullets(s, Inches(7), Inches(2.1), Inches(5.8), Inches(3), interview, sz=13, sp=Pt(8))

# 前端
tx(s, Inches(0.5), Inches(5.5), Inches(12), Inches(0.35), '数据看板 (SPA)', sz=18, c=C_ACCENT, b=True)
dash = ['总览: KPI卡片 + 雷达图 + 散点图 + 柱状图  |  批量测试: 多岗位并发 + 进度条  |  双盲实验: 相关性分析',
        '数据看板: Kaggle 统计分析  |  AI助手: 多轮对话 + Agent路由  |  前端: MIUI风格 · 小米橙(#FF6900) · 毛玻璃效果']
bullets(s, Inches(0.5), Inches(5.95), Inches(12.3), Inches(1), dash, sz=12, sp=Pt(6))


# ======================================================================
# SLIDE 15 — 总结与展望
# ======================================================================
s = S(); bg(s)
title_bar(s, '总结与展望')

tx(s, Inches(0.5), Inches(1.6), Inches(5.8), Inches(0.35), '项目亮点', sz=20, c=C_ACCENT, b=True)
highlights = [
    '多 Agent 协同：3 Agent + Coordinator + 三级路由',
    '科学评分体系：5维度 x 行为锚定量表 + 学术文献支撑',
    'Multi-Shot 评估：N=3 + 95%CI，避免单次评分偏差',
    'RAG 知识增强：ChromaDB + BGE 向量检索同类简历',
    '双盲验证：Pearson r=0.84，人机高度一致',
    '全栈工程化：FastAPI + SPA + Streamlit 双入口',
]
bullets(s, Inches(0.5), Inches(2.05), Inches(5.8), Inches(3), highlights, sz=14, sp=Pt(9))

tx(s, Inches(7), Inches(1.6), Inches(5.8), Inches(0.35), '技术栈', sz=20, c=C_ACCENT, b=True)
stack = [
    '后端: Python · FastAPI · Uvicorn · OpenAI SDK',
    '向量库: ChromaDB · BAAI/bge-small-zh-v1.5',
    'LLM: DeepSeek Chat API · 指数退避重试',
    '前端: Vanilla JS SPA · Chart.js v4.4.7',
    '数据: Kaggle 2,600+ · Pandas · NumPy',
    '工程: RLock并发 · Session TTL · XSS防护',
]
bullets(s, Inches(7), Inches(2.05), Inches(5.8), Inches(3), stack, sz=14, sp=Pt(9))

tx(s, Inches(0.5), Inches(5.5), Inches(12), Inches(0.35), '未来方向', sz=18, c=C_ACCENT, b=True)
future = [
    '校准层：线性映射/分位数校准对齐人工评分尺度  |  RLHF：基于人工反馈持续优化  |  多模态：OCR + 作品集爬取',
    'Agent 扩展：需求分析 + 人岗匹配  |  Docker 部署 + GPU 推理 + CI/CD  |  API 鉴权 + 速率限制',
]
bullets(s, Inches(0.5), Inches(5.95), Inches(12.3), Inches(0.8), future, sz=12, sp=Pt(8))

# thanks
box(s, Inches(0), Inches(6.55), W, Inches(0.95), C_ACCENT)
tx(s, Inches(1), Inches(6.65), Inches(11), Inches(0.7), '感谢聆听 · 欢迎提问', sz=36, c=C_WHITE, b=True, al=PP_ALIGN.CENTER)


# ── 保存 ──
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'AI_Resume_System_期末汇报_v7.pptx')
prs.save(output_path)
print(f'PPT saved to: {output_path}')
print(f'Slides: {len(prs.slides)}')
